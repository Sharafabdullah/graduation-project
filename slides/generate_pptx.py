import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Configuration
SLIDES_JS_PATH = r"d:\University\Graduation Project\slides\slides_data.js"
OUTPUT_PPTX_PATH = r"d:\University\Graduation Project\slides\Graduation_Project_Styled.pptx"
ASSETS_DIR = r"d:\University\Graduation Project\slides"

# Styles from CSS
BG_COLOR = RGBColor(15, 23, 42)      # --bg-dark
TEXT_PRIMARY = RGBColor(248, 250, 252) # --text-primary
TEXT_SECONDARY = RGBColor(148, 163, 184) # --text-secondary
ACCENT_COLOR = RGBColor(56, 189, 248) # --accent
GLASS_BG_COLOR = RGBColor(30, 41, 59) # --glass_bg base
GLASS_BORDER_COLOR = RGBColor(255, 255, 255)

def load_data():
    with open(SLIDES_JS_PATH, 'r', encoding='utf-8') as f:
        content = f.read()
        json_str = re.sub(r'^const\s+slidesData\s*=\s*', '', content).strip().rstrip(';')
        return json.loads(json_str)

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_glass_card(slide, left, top, width, height):
    """Creates a rounded rectangle with semi-transparent dark fill and light border"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = GLASS_BG_COLOR
    fill.transparency = 0.1 # Slight transparency, PPTX doesn't do blur
    
    line = shape.line
    line.color.rgb = GLASS_BORDER_COLOR
    line.width = Pt(1)
    # line.transparency = 0.8 # Line transparency not always supported easily
    
    # Remove text frame auto-resize
    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    return shape

def style_text(paragraph, text, size=Pt(18), color=TEXT_PRIMARY, bold=False, align=None):
    paragraph.text = text
    paragraph.font.size = size
    paragraph.font.color.rgb = color
    paragraph.font.name = 'Arial' # Safe font
    paragraph.font.bold = bold
    if align:
        paragraph.alignment = align

def clean_text(text):
    return re.sub('<[^<]+?>', '', text).strip()

def create_presentation(data):
    prs = Presentation()
    # Use 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank
    
    for slide_data in data:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide)
        s_type = slide_data.get('type')
        
        # --- TITLE SLIDE ---
        if s_type == 'title':
            # Logo
            if 'logo' in slide_data:
                img_path = os.path.join(ASSETS_DIR, slide_data['logo'].replace('assets/', 'assets\\'))
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(5.6), Inches(1), height=Inches(1.5))
            
            # Title Group
            top_start = Inches(3)
            tb = slide.shapes.add_textbox(Inches(1), top_start, Inches(11.3), Inches(2))
            p = tb.text_frame.add_paragraph()
            style_text(p, slide_data.get('subtitle', ''), Pt(24), ACCENT_COLOR, align=PP_ALIGN.CENTER)
            
            p = tb.text_frame.add_paragraph()
            style_text(p, slide_data.get('title', ''), Pt(60), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
            
            # Authors
            if 'authors' in slide_data:
                atb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(1.5))
                p = atb.text_frame.add_paragraph()
                style_text(p, " | ".join(slide_data['authors']), Pt(18), TEXT_PRIMARY, align=PP_ALIGN.CENTER)
                
                p = atb.text_frame.add_paragraph()
                style_text(p, slide_data.get('department', ''), Pt(14), TEXT_SECONDARY, align=PP_ALIGN.CENTER)

        # --- SECTION HEADER ---
        elif s_type == 'section_header':
            # Card in center
            card = add_glass_card(slide, Inches(2), Inches(2), Inches(9.33), Inches(3.5))
            tf = card.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = tf.paragraphs[0]
            style_text(p, slide_data.get('title', ''), Pt(50), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
            
            p = tf.add_paragraph()
            p.text = " "  # Spacer
            
            p = tf.add_paragraph()
            style_text(p, slide_data.get('content', ''), Pt(24), TEXT_SECONDARY, align=PP_ALIGN.CENTER)

        # --- TWO COLUMN / LANDSCAPE ---
        elif s_type == 'two_column' or s_type == 'technological_landscape':
            # Main Title
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            style_text(tb.text_frame.paragraphs[0], slide_data.get('title', ''), Pt(36), ACCENT_COLOR, bold=True)
            
            # Left Card
            left_col = slide_data['columns'][0]
            c1 = add_glass_card(slide, Inches(0.5), Inches(2), Inches(6), Inches(4.5))
            tf = c1.text_frame
            p = tf.paragraphs[0]
            style_text(p, left_col.get('title', ''), Pt(24), ACCENT_COLOR, bold=True)
            p = tf.add_paragraph() # Spacer
            p = tf.add_paragraph()
            style_text(p, left_col.get('content', ''), Pt(18), TEXT_PRIMARY)

            # Right Card
            right_col = slide_data['columns'][1]
            c2 = add_glass_card(slide, Inches(6.8), Inches(2), Inches(6), Inches(4.5))
            tf = c2.text_frame
            p = tf.paragraphs[0]
            style_text(p, right_col.get('title', ''), Pt(24), ACCENT_COLOR, bold=True)
            p = tf.add_paragraph() # Spacer
            p = tf.add_paragraph()
            style_text(p, right_col.get('content', ''), Pt(18), TEXT_PRIMARY)

        # --- GRID ICONS ---
        elif s_type == 'grid_icons':
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            style_text(tb.text_frame.paragraphs[0], slide_data.get('title', ''), Pt(36), ACCENT_COLOR, bold=True)
            
            items = slide_data.get('items', [])
            # 2x2 Grid roughly
            start_x = Inches(0.5)
            start_y = Inches(1.8)
            col_width = Inches(3)
            row_height = Inches(4) 
            gap = Inches(0.2)
            
            # Dynamic grid
            cols = 4 if len(items) >= 4 else len(items)
            w = (12.33 - (gap * (cols-1))) / cols
            
            for i, item in enumerate(items):
                x = start_x + (i * (w + gap))
                card = add_glass_card(slide, x, start_y, w, Inches(4))
                tf = card.text_frame
                p = tf.paragraphs[0]
                # Icon simulation (text)
                style_text(p, item.get('title'), Pt(20), ACCENT_COLOR, bold=True, align=PP_ALIGN.CENTER)
                p = tf.add_paragraph()
                style_text(p, item.get('content'), Pt(16), TEXT_PRIMARY, align=PP_ALIGN.CENTER)

        # --- TWO COLUMN IMAGE ---
        elif s_type == 'two_column_image' or s_type == 'embedded_firmware_logic' or s_type == 'electrical_systems':
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            style_text(tb.text_frame.paragraphs[0], slide_data.get('title', ''), Pt(36), ACCENT_COLOR, bold=True)
            
            if 'subtitle' in slide_data:
                sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
                style_text(sub.text_frame.paragraphs[0], slide_data['subtitle'], Pt(18), TEXT_SECONDARY)

            # Left Content Card
            card = add_glass_card(slide, Inches(0.5), Inches(2), Inches(6), Inches(4.5))
            tf = card.text_frame
            p = tf.paragraphs[0]
            p.text = "" 
            
            if 'bullets' in slide_data:
                for b in slide_data['bullets']:
                    p = tf.add_paragraph()
                    style_text(p, "• " + clean_text(b), Pt(18), TEXT_PRIMARY)
                    p.space_after = Pt(10)

            # Right Image
            img_path = os.path.join(ASSETS_DIR, slide_data.get('image', '').replace('assets/', 'assets\\'))
            if os.path.exists(img_path):
                # Container for image style
                slide.shapes.add_picture(img_path, Inches(6.8), Inches(2), width=Inches(6))

        # --- FEATURE LIST ---
        elif s_type == 'feature_list':
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            style_text(tb.text_frame.paragraphs[0], slide_data.get('title', ''), Pt(36), ACCENT_COLOR, bold=True)
            
            desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1))
            style_text(desc.text_frame.paragraphs[0], slide_data.get('description', ''), Pt(18), TEXT_SECONDARY)
            
            # List items in grid
            feats = slide_data.get('features', [])
            start_y = Inches(2.2)
            
            # 2 columns of features
            mid_point = (len(feats) + 1) // 2
            col1 = feats[:mid_point]
            col2 = feats[mid_point:]
            
            def add_feat_col(items, x_pos):
                curr_y = start_y
                for f in items:
                    # Number circle simulation (just bold text)
                    tb_num = slide.shapes.add_textbox(x_pos, curr_y, Inches(0.5), Inches(0.5))
                    style_text(tb_num.text_frame.paragraphs[0], f['num'], Pt(20), BG_COLOR, bold=True, align=PP_ALIGN.CENTER)
                    # Colored background for number
                    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos, curr_y, Inches(0.5), Inches(0.5))
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = ACCENT_COLOR
                    # Move text explicitly to front? PPTX add order matters. 
                    # Actually, easier to use a Shape with text
                    
                    # Text
                    tb_txt = slide.shapes.add_textbox(x_pos + Inches(0.7), curr_y, Inches(5), Inches(1))
                    tf = tb_txt.text_frame
                    p = tf.paragraphs[0]
                    style_text(p, f['title'], Pt(18), ACCENT_COLOR, bold=True)
                    p = tf.add_paragraph()
                    style_text(p, f['text'], Pt(14), TEXT_PRIMARY)
                    
                    curr_y += Inches(1.2)

            add_feat_col(col1, Inches(0.5))
            add_feat_col(col2, Inches(6.8))

        # --- BLEED IMAGE ---
        elif s_type == 'bleed_image':
             # Split 50/50
             # Left Content
             card = add_glass_card(slide, Inches(0.5), Inches(1), Inches(6), Inches(5.5))
             tf = card.text_frame
             p = tf.paragraphs[0]
             style_text(p, slide_data.get('title', ''), Pt(36), ACCENT_COLOR, bold=True)
             
             for b in slide_data.get('content_blocks', []):
                 p = tf.add_paragraph()
                 p.text = " "
                 p = tf.add_paragraph()
                 style_text(p, b['heading'], Pt(18), ACCENT_COLOR, bold=True)
                 for li in b['list']:
                     p = tf.add_paragraph()
                     style_text(p, "• " + li, Pt(14), TEXT_PRIMARY)
             
             # Right Image
             img_path = os.path.join(ASSETS_DIR, slide_data.get('image', '').replace('assets/', 'assets\\'))
             if os.path.exists(img_path):
                 slide.shapes.add_picture(img_path, Inches(7), Inches(0), width=Inches(6.33), height=Inches(7.5))

        # --- IMAGE GALLERY ---
        elif s_type == 'image_gallery':
             tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
             style_text(tb.text_frame.paragraphs[0], slide_data.get('title', ''), Pt(36), ACCENT_COLOR, bold=True)
             
             items = slide_data.get('items', [])
             cnt = len(items)
             if cnt > 0:
                 w = Inches(3.5)
                 h = Inches(4)
                 margin = Inches(0.5)
                 
                 # Row of cards
                 start_x = (13.333 - (cnt * w) - ((cnt-1)*margin)) / 2
                 
                 for i, item in enumerate(items):
                     x = start_x + (i * (w + margin))
                     # Card background
                     card = add_glass_card(slide, x, Inches(2), w, h)
                     
                     # Image
                     img_path = os.path.join(ASSETS_DIR, item.get('image', '').replace('assets/', 'assets\\'))
                     if os.path.exists(img_path):
                         # Inset image
                         pic = slide.shapes.add_picture(img_path, x + Inches(0.25), Inches(2.25), width=w-Inches(0.5))
                     
                     # Caption
                     tb = slide.shapes.add_textbox(x, Inches(2) + h - Inches(1.5), w, Inches(1))
                     style_text(tb.text_frame.paragraphs[0], item['title'], Pt(18), TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)


        # --- FALLBACK ---
        else:
            tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(2))
            style_text(tb.text_frame.paragraphs[0], slide_data.get('title', 'Unknown'), Pt(36), TEXT_PRIMARY)
    
    prs.save(OUTPUT_PPTX_PATH)
    print(f"Presentation saved to {OUTPUT_PPTX_PATH}")

if __name__ == "__main__":
    data = load_data()
    create_presentation(data)
